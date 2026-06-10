import os
import json
import sys
from pptx import Presentation

sys.path.append(r"C:\Work\PPT Generation application")

from backend.pipeline.queue_manager import STORAGE_DIR

# We want to check:
# 1. Digestive System (proof_batch slide 0)
# 2. Plant Cell (proof_batch slide 1)
# 3. Gravity Diagram (dd6bd672d8ec4e9da5d36a332d22e71d slide 0)
# 4. Water Cycle (proof_batch slide 3)
# 5. Flowchart (proof_batch slide 5)
# 6. Infographic (proof_batch slide 2? Wait, let's see which slide is which)
# 7. Textbook Scan (proof_batch slide 6)

slides_to_report = [
    ("Digestive System", "proof_batch", 0),
    ("Plant Cell", "proof_batch", 1),
    ("Gravity Diagram", "dd6bd672d8ec4e9da5d36a332d22e71d", 0),
    ("Water Cycle", "proof_batch", 3),
    ("Flowchart", "proof_batch", 5),
    ("Infographic", "proof_batch", 2), # Let's verify if slide 2 is solar system/infographic
    ("Textbook Scan", "proof_batch", 6)
]

for name, batch_id, slide_idx in slides_to_report:
    batch_dir = os.path.join(STORAGE_DIR, "batches", batch_id)
    slide_dir = os.path.join(batch_dir, "slides", f"slide_{slide_idx}")
    meta_file = os.path.join(slide_dir, "metadata", "metadata.json")
    
    if not os.path.exists(meta_file):
        print(f"{name}: Metadata file not found at {meta_file}")
        continue
        
    with open(meta_file, "r", encoding="utf-8") as f:
        meta = json.load(f)
        
    components = meta.get("components", [])
    total_extracted = len(components)
    
    useful = 0
    failed = 0
    for c in components:
        # Check if the component has a mask/crop or text, and is visible
        if c.get("visible", True):
            if c.get("type") == "text_label" and c.get("text"):
                useful += 1
            elif c.get("type") in ["image_object", "arrow"] and c.get("mask_path"):
                mask_full_path = os.path.join(slide_dir, c["mask_path"])
                if os.path.exists(mask_full_path):
                    useful += 1
                else:
                    failed += 1
            else:
                failed += 1
        else:
            failed += 1
            
    # Check PPT shape count
    if batch_id == "proof_batch":
        pptx_path = r"C:\Work\PPT Generation application\backend\dry_run_combined_proof.pptx"
    else:
        pptx_path = os.path.join(batch_dir, "presentation.pptx")
        
    shape_count = 0
    if os.path.exists(pptx_path):
        try:
            prs = Presentation(pptx_path)
            # Find slide by index
            if slide_idx < len(prs.slides):
                slide = prs.slides[slide_idx]
                shape_count = len(slide.shapes)
        except Exception as e:
            print(f"Error reading PPTX shape count for {name}: {e}")
            
    print(f"\nReport for: {name}")
    print(f"  Batch ID: {batch_id}, Slide Index: {slide_idx}")
    print(f"  Total Extracted Components: {total_extracted}")
    print(f"  Useful Components: {useful}")
    print(f"  Failed Components: {failed}")
    print(f"  PPT Object Count: {shape_count} (including background image)")
