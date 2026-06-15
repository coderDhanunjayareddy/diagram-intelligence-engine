import os
import cv2
import numpy as np
from typing import Dict, List, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from backend.pipeline.interfaces import BatchJob, SlideMetadata
from backend.pipeline.shape_classifier import ShapeGeometryClassifier, TextStyleExtractor

class PPTGenerator:
    def __init__(self):
        self.shape_classifier = ShapeGeometryClassifier()
        self.text_extractor = TextStyleExtractor()

    def generate_batch_pptx(self, job: BatchJob, batch_dir: str, output_pptx_path: str) -> Dict[int, List[Dict[str, Any]]]:
        """
        Compiles all processed slides in a BatchJob into a single PowerPoint (.pptx) file.
        Replaces custom segmented PNG crops with native PowerPoint vector shapes (rectangles, ovals, diamonds)
        and custom text boxes using programmatically extracted colors, thicknesses, font sizes, and styles.
        Organic illustrations remain transparent PNGs.
        Returns a dictionary mapping slide index to component shape audits:
        {
            slide_idx: [
                {"component_id": str, "predicted_shape": str, "confidence": float, "export_strategy": str},
                ...
            ]
        }
        """
        prs = Presentation()
        
        # Configure slide dimensions dynamically based on first slide aspect ratio
        if job.slides:
            first_slide = job.slides[0]
            img_w = max(1, first_slide.width)
            img_h = max(1, first_slide.height)
            aspect_ratio = img_w / img_h
            if abs(aspect_ratio - 1.0) < 0.05: # Square image (1:1 ratio)
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(10)
            elif aspect_ratio > 1.4: # Widescreen 16:9
                prs.slide_width = Inches(13.33)
                prs.slide_height = Inches(7.5)
            elif aspect_ratio > 1.2: # Standard 4:3
                prs.slide_width = Inches(10)
                prs.slide_height = Inches(7.5)
            else: # Portrait or other custom shape
                prs.slide_height = Inches(10)
                prs.slide_width = Inches(10 * aspect_ratio)
        else:
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
        blank_slide_layout = prs.slide_layouts[6] # blank layout
        
        # Slide shape audits accumulator
        shape_audits = {}
        
        for slide_meta in job.slides:
            slide_idx = slide_meta.slide_index
            slide_audits = []
            
            # 1. Add slide
            slide = prs.slides.add_slide(blank_slide_layout)
            
            # Find the original image path for this slide
            slide_dir = os.path.join(batch_dir, "slides", f"slide_{slide_meta.slide_index}")
            
            # Check if background template exists, otherwise fall back to original
            background_image_path = os.path.join(slide_dir, "masks", "background.png")
            if not os.path.exists(background_image_path):
                background_image_path = os.path.join(slide_dir, "background.png")
                
            if not os.path.exists(background_image_path):
                original_dir = os.path.join(slide_dir, "original")
                if os.path.exists(original_dir) and os.listdir(original_dir):
                    background_image_path = os.path.join(original_dir, os.listdir(original_dir)[0])
                else:
                    original_file = None
                    if os.path.exists(slide_dir):
                        for f in os.listdir(slide_dir):
                            if f.startswith("original."):
                                original_file = f
                                break
                    if original_file:
                        background_image_path = os.path.join(slide_dir, original_file)
                    else:
                        continue

            # Resolve original image path for styling & text extraction
            original_img_path = os.path.join(slide_dir, "original", slide_meta.original_filename)
            if not os.path.exists(original_img_path):
                original_dir = os.path.join(slide_dir, "original")
                if os.path.exists(original_dir) and os.listdir(original_dir):
                    original_img_path = os.path.join(original_dir, os.listdir(original_dir)[0])
                else:
                    original_file = None
                    if os.path.exists(slide_dir):
                        for f in os.listdir(slide_dir):
                            if f.startswith("original."):
                                original_file = f
                                break
                    if original_file:
                        original_img_path = os.path.join(slide_dir, original_file)
            
            # 2. Add Layer 1: Bottom Clean Background Image
            slide.shapes.add_picture(
                background_image_path, 
                Inches(0), Inches(0), 
                prs.slide_width, prs.slide_height
            )
            
            # Coordinate scaling factors: image pixels to Inches
            img_w = max(1, slide_meta.width)
            img_h = max(1, slide_meta.height)
            
            scale_x = float(prs.slide_width) / img_w
            scale_y = float(prs.slide_height) / img_h
            
            # 3. Add Layer 2 and 3: Extracted transparent components and text labels
            sorted_components = sorted(slide_meta.components, key=lambda c: c.z_index)
            
            # Slide-level object counts
            slide_native_count = 0
            slide_png_count = 0
            slide_anim_count = 0
            
            for comp in sorted_components:
                if not comp.visible:
                    continue
                    
                # Compute scaling coordinates (already in EMUs)
                left = int(comp.box[0] * scale_x)
                top = int(comp.box[1] * scale_y)
                width = int(comp.box[2] * scale_x)
                height = int(comp.box[3] * scale_y)
                
                # Safeguard dimensions (avoid zero/negative sizes)
                if width <= 0 or height <= 0:
                    continue
                    
                mask_rel_path = comp.amodal_mask_path or comp.mask_path
                mask_abs_path = ""
                if mask_rel_path:
                    mask_abs_path = os.path.abspath(os.path.join(slide_dir, mask_rel_path))
                    if not os.path.exists(mask_abs_path):
                        # Fallback: check inside the masks/ subdirectory
                        mask_abs_path = os.path.abspath(os.path.join(slide_dir, "masks", os.path.basename(mask_rel_path)))
                
                # Classify the shape geometry and styling
                style_info = self.shape_classifier.classify_shape_and_style(
                    original_img_path=original_img_path,
                    mask_png_path=mask_abs_path,
                    box=comp.box,
                    comp_type=comp.type,
                    semantic_name=comp.semantic_name or ""
                )
                
                # Add to audit report list
                slide_audits.append({
                    "component_id": comp.id,
                    "predicted_shape": style_info["shape_type"],
                    "confidence": style_info["confidence"],
                    "export_strategy": style_info["export_strategy"]
                })
                
                # 3a. Draw Native PowerPoint Shapes (Category A)
                if style_info["export_strategy"] == "Native":
                    shape_type = style_info["shape_type"]
                    fill_rgb = style_info["fill_color"]
                    border_rgb = style_info["border_color"]
                    border_thickness = style_info["border_thickness"]
                    corner_radius = style_info["corner_radius"]
                    
                    try:
                        # Map to appropriate PPT Shape Enum
                        ppt_shape_enum = MSO_SHAPE.RECTANGLE
                        if shape_type == "Rectangle":
                            ppt_shape_enum = MSO_SHAPE.RECTANGLE
                        elif shape_type == "Rounded Rectangle":
                            ppt_shape_enum = MSO_SHAPE.ROUNDED_RECTANGLE
                        elif shape_type == "Circle":
                            ppt_shape_enum = MSO_SHAPE.OVAL
                        elif shape_type == "Diamond":
                            ppt_shape_enum = MSO_SHAPE.DIAMOND
                        elif shape_type in ["Line", "Arrow"]:
                            # Draw thin line block to ensure rotation/placement correctness
                            ppt_shape_enum = MSO_SHAPE.RECTANGLE
                            
                        shape = slide.shapes.add_shape(ppt_shape_enum, left, top, width, height)
                        
                        # Style solid fill
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
                        
                        # Style border line
                        shape.line.color.rgb = RGBColor(*border_rgb)
                        shape.line.width = Pt(border_thickness)
                        
                        # Style corner radius adjustment for rounded rectangles
                        if shape_type == "Rounded Rectangle" and corner_radius > 0:
                            shape.adjustments[0] = corner_radius
                            
                        # If Arrow, draw arrowhead using simulated polygon on ends or properties
                        if shape_type == "Arrow":
                            # Note: line.end_arrowhead is only supported on connector shapes in python-pptx,
                            # but we can set line formatting or standard shape styling here if supported.
                            pass
                            
                        slide_native_count += 1
                        slide_anim_count += 1
                        
                    except Exception as e_shape:
                        print(f"Error rendering native PPT shape {comp.id}: {e_shape}")
                        # Fallback to PNG
                        if mask_abs_path and os.path.exists(mask_abs_path):
                            try:
                                slide.shapes.add_picture(mask_abs_path, left, top, width, height)
                                slide_png_count += 1
                                slide_anim_count += 1
                            except Exception:
                                pass
                                
                # 3b. Text Labels (using standard TextBoxes with extracted font styling)
                elif comp.type == "text_label" and comp.text:
                    try:
                        # Extract font styling parameters
                        text_style = self.text_extractor.extract_text_style(
                            original_img_path=original_img_path,
                            box=comp.box,
                            text_content=comp.text
                        )
                        
                        # Add native TextBox shape
                        shape = slide.shapes.add_textbox(left, top, width, height)
                        
                        # Style background and line border transparent
                        shape.fill.background()
                        shape.line.fill.background()
                        
                        tf = shape.text_frame
                        tf.word_wrap = True
                        tf.margin_left = Inches(0.04)
                        tf.margin_right = Inches(0.04)
                        tf.margin_top = Inches(0.02)
                        tf.margin_bottom = Inches(0.02)
                        
                        lines = comp.text.split("\n")
                        for p_idx, line in enumerate(lines):
                            if p_idx == 0:
                                p = tf.paragraphs[0]
                            else:
                                p = tf.add_paragraph()
                                
                            p.text = line
                            p.font.name = 'Arial'
                            p.font.size = Pt(text_style["estimated_font_size"])
                            p.font.bold = text_style["bold_estimate"]
                            p.font.color.rgb = RGBColor(*text_style["font_color"])
                            
                            # Alignment mapping
                            align_str = text_style["alignment"]
                            if align_str == "left":
                                p.alignment = PP_ALIGN.LEFT
                            elif align_str == "right":
                                p.alignment = PP_ALIGN.RIGHT
                            else:
                                p.alignment = PP_ALIGN.CENTER
                                
                        slide_native_count += 1
                        slide_anim_count += 1
                        
                    except Exception as e_text:
                        print(f"Error rendering native PPT textbox {comp.id}: {e_text}")
                        
                # 3c. Segmented Transparent PNG Layers (Category C Organic Visual Elements)
                else:
                    if mask_abs_path and os.path.exists(mask_abs_path):
                        try:
                            slide.shapes.add_picture(mask_abs_path, left, top, width, height)
                            slide_png_count += 1
                            slide_anim_count += 1
                        except Exception as e_pic:
                            print(f"Error embedding transparent crop {comp.id}: {e_pic}")
                            
            # Store counts in slide metadata or dictionary for reference
            shape_audits[slide_meta.slide_index] = slide_audits
            slide_meta.routing_status = f"Native shapes: {slide_native_count}, PNG: {slide_png_count}"
            
        # Save presentation to output path
        prs.save(output_pptx_path)
        print(f"PowerPoint Presentation exported to: {output_pptx_path}")
        return shape_audits
